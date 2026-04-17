"""Document parser: extracts clean text from PDF, DOCX, PPTX, TXT files."""

import re
from pathlib import Path

import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}


def extract_text(filepath: str) -> str:
    """Extract text from a document file and return cleaned text."""
    path = Path(filepath)
    ext = path.suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported file format: {ext}. Supported: {', '.join(SUPPORTED_EXTENSIONS)}")

    if ext == ".pdf":
        text = _extract_pdf(path)
    elif ext == ".docx":
        text = _extract_docx(path)
    elif ext == ".pptx":
        text = _extract_pptx(path)
    elif ext == ".txt":
        text = _extract_txt(path)

    return _clean_text(text)


def _extract_pdf(path: Path) -> str:
    doc = fitz.open(str(path))
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return "\n".join(pages)


def _extract_docx(path: Path) -> str:
    doc = DocxDocument(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def _extract_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    slides_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slides_text.append(text)
    return "\n".join(slides_text)


def _extract_txt(path: Path) -> str:
    encodings = ["utf-8", "cp1251", "latin-1"]
    for enc in encodings:
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, ValueError):
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def _clean_text(text: str) -> str:
    """Clean extracted text: remove page numbers, headers/footers, normalize whitespace."""
    # Remove standalone page numbers (e.g. lines that are just "12" or "- 12 -")
    text = re.sub(r"^\s*[-–—]?\s*\d{1,4}\s*[-–—]?\s*$", "", text, flags=re.MULTILINE)

    # Collapse multiple blank lines into one
    text = re.sub(r"\n{3,}", "\n\n", text)

    # Normalize spaces (but keep newlines)
    text = re.sub(r"[^\S\n]+", " ", text)

    # Remove leading/trailing whitespace per line
    lines = [line.strip() for line in text.split("\n")]
    text = "\n".join(lines)

    return text.strip()
